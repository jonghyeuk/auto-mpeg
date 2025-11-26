"""
PPT 요소 추출기
PPT에서 텍스트, 이미지, 도형을 개별적으로 추출
"""
from pathlib import Path
from typing import Dict, List, Any
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io


def extract_ppt_elements(pptx_path: Path, output_json: Path, output_dir: Path) -> Dict[str, Any]:
    """
    PPT 파일에서 요소들을 추출

    Args:
        pptx_path: PPT 파일 경로
        output_json: 출력 JSON 파일 경로
        output_dir: 추출된 이미지 저장 디렉토리

    Returns:
        추출된 요소 정보
    """
    output_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(pptx_path))

    slides_data = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        print(f"\n📄 슬라이드 {slide_idx} 처리 중...")

        slide_data = {
            "index": slide_idx,
            "texts": [],
            "images": []
        }

        # 슬라이드 내 모든 shape 순회
        for shape_idx, shape in enumerate(slide.shapes):

            # 텍스트 추출
            if hasattr(shape, "text") and shape.text.strip():
                text_info = {
                    "text": shape.text.strip(),
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height
                }
                slide_data["texts"].append(text_info)
                print(f"  📝 텍스트: {shape.text.strip()[:50]}...")

            # 이미지 추출
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    image_bytes = image.blob

                    # 이미지 파일로 저장
                    image_filename = f"slide_{slide_idx:03d}_img_{shape_idx}.png"
                    image_path = output_dir / image_filename

                    # PIL로 이미지 처리 및 저장
                    pil_image = Image.open(io.BytesIO(image_bytes))
                    pil_image.save(str(image_path), "PNG")

                    image_info = {
                        "path": str(image_path.relative_to(output_dir.parent)),
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height
                    }
                    slide_data["images"].append(image_info)
                    print(f"  🖼️  이미지: {image_filename}")

                except Exception as e:
                    print(f"  ⚠️  이미지 추출 실패: {e}")

        slides_data.append(slide_data)

    # JSON으로 저장
    output_data = {
        "slides": slides_data,
        "total_slides": len(slides_data)
    }

    output_json.parent.mkdir(parents=True, exist_ok=True)
    with open(output_json, 'w', encoding='utf-8') as f:
        json.dump(output_data, f, ensure_ascii=False, indent=2)

    print(f"\n✅ 요소 추출 완료:")
    print(f"  - 슬라이드 수: {len(slides_data)}")
    print(f"  - JSON: {output_json}")
    print(f"  - 이미지: {output_dir}")

    return output_data


if __name__ == "__main__":
    # 테스트 코드
    import sys

    if len(sys.argv) < 2:
        print("Usage: python ppt_element_extractor.py <pptx_file>")
        sys.exit(1)

    pptx_file = Path(sys.argv[1])
    project_root = Path(__file__).parent.parent.parent
    output_json = project_root / "data" / "temp" / "reactant" / "elements.json"
    output_imgs = project_root / "data" / "temp" / "reactant" / "elements"

    extract_ppt_elements(pptx_file, output_json, output_imgs)
