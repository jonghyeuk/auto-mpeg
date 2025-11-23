"""
모듈 A: PPT 파서
PPT 파일에서 슬라이드 텍스트, 이미지 정보를 추출하여 JSON으로 변환
"""
import json
from pathlib import Path
from typing import List, Dict, Any
from pptx import Presentation
from pptx.util import Inches
import io
from PIL import Image


class PPTParser:
    """PPT 파일을 파싱하여 슬라이드 정보를 추출하는 클래스"""

    def __init__(self, ppt_path: str):
        """
        Args:
            ppt_path: PPT 파일 경로
        """
        self.ppt_path = Path(ppt_path)
        self.presentation = Presentation(str(self.ppt_path))

    def extract_text_from_shape(self, shape) -> str:
        """슬라이드 Shape에서 텍스트 추출"""
        if hasattr(shape, "text"):
            return shape.text.strip()
        return ""

    def extract_slide_text(self, slide) -> Dict[str, str]:
        """슬라이드에서 제목, 본문, 노트 텍스트 추출"""
        title = ""
        body_parts = []

        # 슬라이드 내 모든 Shape에서 텍스트 추출
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    # 첫 번째 큰 텍스트를 제목으로 간주
                    if not title and len(text) < 100:
                        title = text
                    else:
                        body_parts.append(text)

        # 노트 추출
        notes = ""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                notes = notes_slide.notes_text_frame.text.strip()

        return {
            "title": title,
            "body": "\n".join(body_parts),
            "notes": notes
        }

    def save_slide_as_image(self, slide_index: int, output_path: Path) -> None:
        """
        슬라이드를 이미지로 저장
        Note: python-pptx는 직접 이미지 저장을 지원하지 않으므로,
        실제로는 LibreOffice/PowerPoint CLI 또는 pptx2png 같은 도구가 필요합니다.
        여기서는 placeholder로 구현합니다.
        """
        # TODO: LibreOffice headless mode 또는 다른 변환 도구 사용
        # 예: libreoffice --headless --convert-to png --outdir output_dir input.pptx
        pass

    def parse(self, output_json_path: Path, output_img_dir: Path) -> List[Dict[str, Any]]:
        """
        PPT를 파싱하여 슬라이드 정보를 JSON으로 저장하고 이미지 추출

        Args:
            output_json_path: 출력 JSON 파일 경로
            output_img_dir: 슬라이드 이미지 저장 디렉토리

        Returns:
            슬라이드 정보 리스트
        """
        output_img_dir.mkdir(parents=True, exist_ok=True)

        slides_data = []

        for idx, slide in enumerate(self.presentation.slides, start=1):
            # 텍스트 정보 추출
            text_data = self.extract_slide_text(slide)

            # 이미지 파일명
            img_filename = f"slide_{idx:03d}.png"
            img_path = output_img_dir / img_filename

            slide_info = {
                "index": idx,
                "title": text_data["title"],
                "body": text_data["body"],
                "notes": text_data["notes"],
                "image": str(img_path.relative_to(output_json_path.parent.parent))
            }

            slides_data.append(slide_info)

            # 슬라이드 이미지 저장 (실제 구현 필요)
            # self.save_slide_as_image(idx, img_path)

        # JSON 저장
        output_json_path.parent.mkdir(parents=True, exist_ok=True)
        with open(output_json_path, 'w', encoding='utf-8') as f:
            json.dump(slides_data, f, ensure_ascii=False, indent=2)

        print(f"✓ PPT 파싱 완료: {len(slides_data)}개 슬라이드")
        print(f"  - JSON: {output_json_path}")
        print(f"  - 이미지: {output_img_dir}")

        return slides_data


def find_libreoffice_path():
    """
    시스템에서 LibreOffice 실행 파일 경로 찾기

    Returns:
        LibreOffice 실행 파일 경로 또는 None
    """
    import os
    import platform

    system = platform.system()

    if system == "Windows":
        # Windows에서 가능한 경로들
        possible_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            os.path.expandvars(r"%ProgramFiles%\LibreOffice\program\soffice.exe"),
            os.path.expandvars(r"%ProgramFiles(x86)%\LibreOffice\program\soffice.exe"),
        ]

        for path in possible_paths:
            if Path(path).exists():
                return path

        return None

    elif system == "Darwin":  # macOS
        return "/Applications/LibreOffice.app/Contents/MacOS/soffice"

    else:  # Linux
        return "libreoffice"


def convert_pptx_to_images(pptx_path: Path, output_dir: Path) -> None:
    """
    LibreOffice를 사용하여 PPTX를 PNG 이미지로 변환

    LibreOffice는 먼저 PDF로 변환하고, 그 다음 각 페이지를 PNG로 변환합니다.

    Args:
        pptx_path: PPTX 파일 경로
        output_dir: 출력 디렉토리
    """
    import subprocess
    import platform
    import time

    output_dir.mkdir(parents=True, exist_ok=True)

    # LibreOffice 실행 파일 찾기
    libreoffice_path = find_libreoffice_path()

    if not libreoffice_path:
        raise FileNotFoundError(
            "LibreOffice를 찾을 수 없습니다. "
            "다음 위치에 설치되어 있는지 확인하세요:\n"
            "  - C:\\Program Files\\LibreOffice\\program\\soffice.exe\n"
            "  - C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe"
        )

    print(f"📄 LibreOffice 경로: {libreoffice_path}")

    # Step 1: PPTX → PDF 변환
    print(f"🔄 STEP 1: PPTX → PDF 변환 중...")
    temp_pdf = output_dir / f"{pptx_path.stem}.pdf"

    cmd_pdf = [
        libreoffice_path,
        "--headless",
        "--convert-to", "pdf",
        "--outdir", str(output_dir),
        str(pptx_path)
    ]

    try:
        result = subprocess.run(cmd_pdf, capture_output=True, text=True, timeout=120)

        if result.returncode != 0:
            print(f"✗ PDF 변환 실패: {result.stderr}")
            raise RuntimeError(f"LibreOffice PDF 변환 실패: {result.stderr}")

        # PDF 파일 생성 대기
        time.sleep(2)

        if not temp_pdf.exists():
            raise FileNotFoundError(f"PDF 파일이 생성되지 않았습니다: {temp_pdf}")

        print(f"✓ PDF 변환 완료: {temp_pdf.name}")

        # Step 2: PDF → PNG 변환 (각 페이지)
        print(f"🔄 STEP 2: PDF → PNG 변환 중...")

        try:
            import fitz  # PyMuPDF

            # PDF 열기
            pdf_document = fitz.open(str(temp_pdf))
            page_count = pdf_document.page_count

            print(f"  - {page_count}개 페이지 발견")

            # 각 페이지를 개별 PNG로 저장 (스케일+패딩 적용)
            import cv2
            import numpy as np

            target_width = 1920
            target_height = 1080

            for page_num in range(page_count):
                page = pdf_document[page_num]

                # 페이지를 이미지로 렌더링 (150 DPI)
                mat = fitz.Matrix(150/72, 150/72)  # 72 DPI -> 150 DPI
                pix = page.get_pixmap(matrix=mat)

                # NumPy 배열로 변환
                img_data = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, pix.n)

                # RGB로 변환 (RGBA인 경우)
                if pix.n == 4:
                    img = cv2.cvtColor(img_data, cv2.COLOR_RGBA2BGR)
                else:
                    img = cv2.cvtColor(img_data, cv2.COLOR_RGB2BGR)

                # 원본 크기
                orig_height, orig_width = img.shape[:2]

                # 비율 유지하면서 목표 크기에 맞추기
                scale = min(target_width / orig_width, target_height / orig_height)
                new_width = int(orig_width * scale)
                new_height = int(orig_height * scale)

                # 리사이즈
                resized = cv2.resize(img, (new_width, new_height), interpolation=cv2.INTER_AREA)

                # 중앙 정렬을 위한 패딩 계산
                pad_x = (target_width - new_width) // 2
                pad_y = (target_height - new_height) // 2

                # 검은색 배경에 이미지 배치 (1920x1080 보장)
                result = np.zeros((target_height, target_width, 3), dtype=np.uint8)
                result[pad_y:pad_y+new_height, pad_x:pad_x+new_width] = resized

                # PNG로 저장
                output_path = output_dir / f"slide_{page_num + 1:03d}.png"
                cv2.imwrite(str(output_path), result)

                print(f"  - 슬라이드 {page_num + 1} 저장: {output_path.name} (원본: {orig_width}x{orig_height} → {target_width}x{target_height})")

            pdf_document.close()

            print(f"✓ PPTX → PNG 변환 완료: {output_dir}")
            print(f"  - 총 {page_count}개 슬라이드 이미지 생성")

            # 임시 PDF 파일 삭제
            temp_pdf.unlink()

        except ImportError:
            print("⚠️  PyMuPDF 모듈이 설치되어 있지 않습니다.")
            print("   pip install pymupdf 를 실행하세요.")
            raise ImportError("PyMuPDF 모듈 필요")

    except FileNotFoundError:
        raise FileNotFoundError(
            "LibreOffice가 설치되어 있지 않습니다.\n"
            "https://www.libreoffice.org/download/download/ 에서 다운로드하세요."
        )
    except subprocess.TimeoutExpired:
        raise TimeoutError("LibreOffice 변환 시간 초과 (120초)")


if __name__ == "__main__":
    # 테스트 코드
    import sys
    from pathlib import Path

    if len(sys.argv) < 2:
        print("Usage: python ppt_parser.py <pptx_file>")
        sys.exit(1)

    pptx_file = Path(sys.argv[1])
    project_root = Path(__file__).parent.parent.parent
    output_json = project_root / "data" / "meta" / "slides.json"
    output_imgs = project_root / "data" / "temp" / "slides_img"

    parser = PPTParser(str(pptx_file))
    slides = parser.parse(output_json, output_imgs)

    # 이미지 변환 (LibreOffice 사용)
    convert_pptx_to_images(pptx_file, output_imgs)
